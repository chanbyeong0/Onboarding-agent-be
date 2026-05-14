"""
test_document 모듈은 문서 업로드 API와 텍스트 청크 분리 로직을 검증한다.
PPTX 파일을 메모리에서 생성해 실제 파서 경로를 테스트한다.
"""

from io import BytesIO

import pytest
from pptx import Presentation

from app.models.document import DocumentModel
from app.services.document_service import chunk_text


def build_pptx_bytes(text: str) -> bytes:
    """테스트용 PPTX 파일 바이트를 생성한다.

    Args:
        text: 슬라이드에 넣을 텍스트.

    Returns:
        bytes: 업로드 요청에 사용할 PPTX 바이너리.
    """

    # python-pptx로 실제 파서가 읽을 수 있는 프레젠테이션을 만든다
    presentation = Presentation()

    # 빈 슬라이드 레이아웃에 텍스트 박스를 추가해 파싱 대상 텍스트를 넣는다
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(0, 0, 5000000, 500000)
    text_box.text = text

    buffer = BytesIO()

    # 프레젠테이션을 메모리 버퍼에 저장해 업로드 파일처럼 사용한다
    presentation.save(buffer)
    return buffer.getvalue()


def test_chunk_text_with_overlap() -> None:
    """chunk_text가 overlap을 유지하며 텍스트를 나누는지 확인한다."""

    # 짧은 테스트 문자열을 작은 청크 크기로 나눠 겹침 동작을 검증한다
    chunks = chunk_text("abcdefghijklmnopqrstuvwxyz", chunk_size=10, overlap=2)

    assert chunks == ["abcdefghij", "ijklmnopqr", "qrstuvwxyz"]


@pytest.mark.asyncio
async def test_upload_pptx_document(client, auth_headers) -> None:
    """PPTX 업로드 API가 문서를 저장하고 응답하는지 확인한다."""

    # 테스트용 PPTX 파일을 메모리에서 생성한다
    pptx_content = build_pptx_bytes("신입사원 온보딩 안내")

    files = {
        "file": (
            "onboarding.pptx",
            pptx_content,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    }

    # 인증 헤더와 multipart 파일로 문서 업로드 API를 호출한다
    response = await client.post("/api/v1/documents/upload", files=files, headers=auth_headers)

    assert response.status_code == 201
    body = response.json()
    assert body["title"] == "onboarding.pptx"
    assert body["file_type"] == "pptx"

    # MongoDB에 문서 메타데이터가 실제로 저장되었는지 확인한다
    documents = await DocumentModel.find_all().to_list()
    assert len(documents) == 1
    assert documents[0].title == "onboarding.pptx"
