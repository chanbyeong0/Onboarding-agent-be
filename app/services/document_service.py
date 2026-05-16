"""
document_service 모듈은 업로드 파일 저장, PDF/PPT 텍스트 추출, 청크 생성을 담당한다.
RAG 색인 호출 지점은 TODO로 남겨 두고, 문서 메타데이터는 MongoDB에 저장한다.
"""

from pathlib import Path
import subprocess
from typing import Any
from uuid import uuid4

import anyio
import pdfplumber
import pypdfium2 as pdfium
from fastapi import HTTPException, UploadFile, status
from pptx import Presentation

from app.core.config import settings
from app.crud import document as document_crud
from app.crud import document_page as document_page_crud
from app.models.document import DocumentModel

SUPPORTED_EXTENSIONS = {"pdf", "ppt", "pptx"}


def get_file_type(filename: str) -> str:
    """파일명에서 지원 가능한 문서 유형을 추출한다.

    Args:
        filename: 업로드된 원본 파일명.

    Returns:
        str: pdf, ppt, pptx 중 하나.

    Raises:
        HTTPException: 지원하지 않는 확장자일 때 발생한다.
    """

    extension = Path(filename).suffix.lower().lstrip(".")
    if extension not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="PDF, PPT, PPTX 파일만 업로드할 수 있습니다.",
        )
    return extension


async def save_upload_file(file: UploadFile) -> Path:
    """업로드 파일을 서버 디스크에 저장한다.

    Args:
        file: FastAPI가 전달한 multipart 업로드 파일.

    Returns:
        Path: 저장된 파일의 경로.
    """

    # 파일 유형 검증으로 허용되지 않은 확장자의 저장을 차단한다
    file_type = get_file_type(file.filename or "")
    upload_dir = Path(settings.upload_dir)

    # 업로드 디렉터리가 없으면 생성해 이후 파일 저장 실패를 방지한다
    await anyio.to_thread.run_sync(lambda: upload_dir.mkdir(parents=True, exist_ok=True))

    safe_name = f"{uuid4().hex}.{file_type}"
    file_path = upload_dir / safe_name

    # UploadFile의 비동기 read로 요청 본문을 메모리에 읽는다
    content = await file.read()

    # 디스크 쓰기는 블로킹 작업이므로 스레드로 넘겨 이벤트 루프를 막지 않는다
    await anyio.to_thread.run_sync(file_path.write_bytes, content)
    return file_path


def extract_pdf_text(file_path: Path) -> list[dict[str, Any]]:
    """PDF 파일에서 페이지별 텍스트를 추출한다.

    Args:
        file_path: 서버에 저장된 PDF 파일 경로.

    Returns:
        list[dict[str, Any]]: page_number와 text를 포함한 페이지 목록.
    """

    pages: list[dict[str, Any]] = []

    # pdfplumber가 PDF 페이지 객체를 열어 텍스트 추출 API를 제공한다
    with pdfplumber.open(file_path) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            # PDF 페이지에서 텍스트를 추출하고 빈 페이지는 빈 문자열로 정규화한다
            text = page.extract_text() or ""
            pages.append({"page_number": index, "text": text})
    return pages


def extract_ppt_text(file_path: Path) -> list[dict[str, Any]]:
    """PPT/PPTX 파일에서 슬라이드별 텍스트를 추출한다.

    Args:
        file_path: 서버에 저장된 PPT/PPTX 파일 경로.

    Returns:
        list[dict[str, Any]]: page_number와 text를 포함한 슬라이드 목록.
    """

    pages: list[dict[str, Any]] = []

    # python-pptx가 프레젠테이션 파일을 슬라이드 단위 객체로 변환한다
    presentation = Presentation(file_path)
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        pages.append({"page_number": index, "text": "\n".join(texts)})
    return pages


async def extract_pages(file_path: Path, file_type: str) -> list[dict[str, Any]]:
    """파일 유형에 맞는 파서를 선택해 페이지 텍스트를 추출한다.

    Args:
        file_path: 서버에 저장된 파일 경로.
        file_type: pdf, ppt, pptx 중 하나.

    Returns:
        list[dict[str, Any]]: 페이지 번호와 텍스트 목록.
    """

    if file_type == "pdf":
        # PDF 파싱은 블로킹 작업이므로 스레드로 넘겨 처리한다
        return await anyio.to_thread.run_sync(extract_pdf_text, file_path)

    # PPT/PPTX 파싱은 블로킹 작업이므로 스레드로 넘겨 처리한다
    return await anyio.to_thread.run_sync(extract_ppt_text, file_path)


def convert_office_to_pdf(file_path: Path) -> Path:
    """LibreOffice headless 모드로 PPT/PPTX를 PDF로 변환한다."""

    output_dir = file_path.parent
    completed = subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(file_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    pdf_path = output_dir / f"{file_path.stem}.pdf"
    if completed.returncode != 0 or not pdf_path.exists():
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"문서 PDF 변환에 실패했습니다: {completed.stderr or completed.stdout}",
        )
    return pdf_path


def render_pdf_pages(pdf_path: Path, document_id: str) -> list[dict[str, Any]]:
    """PDF 페이지를 PNG 이미지로 렌더링한다."""

    page_dir = Path(settings.upload_dir) / "pages" / document_id
    page_dir.mkdir(parents=True, exist_ok=True)
    rendered_pages: list[dict[str, Any]] = []

    pdf = pdfium.PdfDocument(pdf_path)
    try:
        for index in range(len(pdf)):
            page_number = index + 1
            image_path = page_dir / f"{page_number}.png"
            page = pdf[index]
            bitmap = page.render(scale=2).to_pil()
            bitmap.save(image_path)
            rendered_pages.append({"page_number": page_number, "image_path": str(image_path)})
    finally:
        pdf.close()
    return rendered_pages


async def render_document_pages(file_path: Path, file_type: str, document_id: str) -> list[dict[str, Any]]:
    """문서 유형에 맞게 페이지 이미지를 생성한다."""

    source_pdf = await get_viewer_pdf_path(file_path, file_type)
    return await anyio.to_thread.run_sync(render_pdf_pages, source_pdf, document_id)


async def get_viewer_pdf_path(file_path: Path, file_type: str) -> Path:
    """프론트 PDF 뷰어가 사용할 PDF 경로를 반환한다."""

    if file_type == "pdf":
        return file_path
    return await anyio.to_thread.run_sync(convert_office_to_pdf, file_path)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """텍스트를 겹침 구간이 있는 고정 길이 청크로 분리한다.

    Args:
        text: 청크로 나눌 원문 텍스트.
        chunk_size: 하나의 청크에 담을 최대 문자 수.
        overlap: 인접 청크가 공유할 문자 수.

    Returns:
        list[str]: 공백을 정리한 텍스트 청크 목록.
    """

    normalized = " ".join(text.split())
    if not normalized:
        return []
    if overlap >= chunk_size:
        raise ValueError("overlap은 chunk_size보다 작아야 합니다.")

    chunks: list[str] = []
    start = 0
    while start < len(normalized):
        end = start + chunk_size
        chunks.append(normalized[start:end])
        if end >= len(normalized):
            break
        start = end - overlap
    return chunks


def chunk_pages(
    pages: list[dict[str, Any]],
    document_id: str,
    chunk_size: int = 500,
    overlap: int = 50,
) -> list[dict[str, Any]]:
    """페이지 목록을 메타데이터가 포함된 청크 목록으로 변환한다.

    Args:
        pages: page_number와 text를 담은 페이지 목록.
        document_id: 청크가 속한 문서 ObjectId 문자열.
        chunk_size: 하나의 청크에 담을 최대 문자 수.
        overlap: 인접 청크가 공유할 문자 수.

    Returns:
        list[dict[str, Any]]: text, document_id, page_number를 담은 청크 목록.
    """

    chunks: list[dict[str, Any]] = []
    for page in pages:
        # 페이지 텍스트를 RAG 색인에 적합한 작은 단위로 분리한다
        page_chunks = chunk_text(page["text"], chunk_size=chunk_size, overlap=overlap)
        for chunk in page_chunks:
            chunks.append(
                {
                    "text": chunk,
                    "document_id": document_id,
                    "page_number": page["page_number"],
                }
            )
    return chunks


async def process_upload(file: UploadFile) -> DocumentModel:
    """문서 업로드 파일을 저장하고 파싱 가능한 메타데이터를 생성한다.

    Args:
        file: 사용자가 업로드한 PDF/PPT/PPTX 파일.

    Returns:
        DocumentModel: MongoDB에 저장된 문서 메타데이터.
    """

    if not file.filename:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="파일명이 없는 업로드는 처리할 수 없습니다.",
        )

    # 파일 유형을 먼저 검증해 저장과 파싱 흐름을 결정한다
    file_type = get_file_type(file.filename)

    # 업로드 파일을 디스크에 저장하고 저장 경로를 회수한다
    file_path = await save_upload_file(file)

    # PPT/PPTX는 브라우저 뷰어용 PDF로 변환하고, PDF는 원본을 그대로 사용한다
    viewer_pdf_path = await get_viewer_pdf_path(file_path, file_type)

    # MongoDB(Beanie)에 문서 메타데이터를 저장한다
    document = await document_crud.create(
        title=file.filename,
        file_path=str(file_path),
        file_type=file_type,
        viewer_pdf_path=str(viewer_pdf_path),
    )

    # 파일 유형에 맞는 파서로 페이지(슬라이드) 단위 텍스트를 추출한다
    pages = await extract_pages(file_path, file_type)

    # 브라우저 강의 플레이어에서 사용할 페이지 이미지를 생성한다
    rendered_pages = await render_document_pages(file_path, file_type, str(document.id))
    text_by_page = {page["page_number"]: page["text"] for page in pages}
    await document_page_crud.create_many(
        str(document.id),
        [
            {
                "page_number": rendered_page["page_number"],
                "image_path": rendered_page["image_path"],
                "text": text_by_page.get(rendered_page["page_number"], ""),
            }
            for rendered_page in rendered_pages
        ],
    )

    # 페이지 텍스트를 청크 단위로 잘라 메타데이터를 부착한다
    _chunks = chunk_pages(pages, document_id=str(document.id))

    # TODO: 추후 구현 - RAG 파이프라인에 청크를 색인 요청
    # await rag_pipeline.index_chunks(_chunks, metadata={"document_id": str(document.id)})

    return document
